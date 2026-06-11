import json
import os
import re
import shutil
import traceback
import uuid
import zipfile
from datetime import datetime
from io import BytesIO
from pathlib import Path

print("APP.PY LOADED - v1.0")

import fitz
from docx import Document
from docx.shared import Inches, Pt
from flask import Flask, jsonify, render_template, request, send_file
from openpyxl import Workbook
from openpyxl.styles import Font
from werkzeug.exceptions import HTTPException, RequestEntityTooLarge

try:
    from PIL import Image  # noqa: F401
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

BASE_DIR = Path(__file__).resolve().parent
DOWNLOADS_DIR = BASE_DIR / "descargas"
METADATA_FILE = "metadata.json"
SESSION_FILE = "session.json"
ALLOWED_EXPORT_SUFFIXES = {".pdf", ".zip", ".docx", ".xlsx", ".pptx"}
MAX_COMPRESSED_BYTES = 999 * 1024
SOURCE_FILE = "source.bin"

app = Flask(__name__)
app.secret_key = os.urandom(24).hex()
app.config["UPLOAD_FOLDER"] = DOWNLOADS_DIR
app.config["MAX_CONTENT_LENGTH"] = 200 * 1024 * 1024
WEBP_INSERT_SUPPORTED = None

@app.route("/api/test2", methods=["GET", "POST"])
def api_test2():
    print("API_TEST2 CALLED")
    return jsonify({"message": "Test2 OK"})


def ensure_upload_dir():
    DOWNLOADS_DIR.mkdir(parents=True, exist_ok=True)


def supports_webp_insert():
    global WEBP_INSERT_SUPPORTED

    if WEBP_INSERT_SUPPORTED is not None:
        return WEBP_INSERT_SUPPORTED

    if not PIL_AVAILABLE:
        WEBP_INSERT_SUPPORTED = False
        return WEBP_INSERT_SUPPORTED

    try:
        probe_doc = fitz.open()
        probe_page = probe_doc.new_page(width=10, height=10)
        probe_pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 2, 2), False)
        probe_pixmap.clear_with(255)
        probe_bytes = probe_pixmap.pil_tobytes(format="WEBP", quality=60, method=6)
        probe_page.insert_image(fitz.Rect(0, 0, 10, 10), stream=probe_bytes)
        WEBP_INSERT_SUPPORTED = True
    except Exception:
        WEBP_INSERT_SUPPORTED = False
    finally:
        try:
            probe_doc.close()
        except Exception:
            pass

    return WEBP_INSERT_SUPPORTED


def remove_session(session_id):
    session_path = DOWNLOADS_DIR / session_id
    if session_path.exists():
        shutil.rmtree(session_path)


def build_image_metadata(doc):
    metadata_by_xref = {}

    for page_num in range(len(doc)):
        page = doc[page_num]
        for image_index, image in enumerate(page.get_images(full=True)):
            xref = image[0]
            entry = metadata_by_xref.get(xref)
            if entry is None:
                base = doc.extract_image(xref)
                ext = base["ext"]
                filename = f"xref_{xref:05d}.{ext}"
                entry = {
                    "xref": xref,
                    "name": filename,
                    "ext": ext,
                    "width": base["width"],
                    "height": base["height"],
                    "pages": [],
                }
                metadata_by_xref[xref] = entry

            bbox = page.get_image_bbox(image)
            entry["pages"].append({
                "page": page_num,
                "index": image_index,
                "rect": [bbox.x0, bbox.y0, bbox.x1, bbox.y1],
            })

    return list(metadata_by_xref.values())


def save_extracted_images(doc, images_dir, metadata):
    for item in metadata:
        image_data = doc.extract_image(item["xref"])
        output_path = images_dir / item["name"]
        ext = item["ext"].lower()

        if ext in {"jpg", "jpeg"}:
            output_path.write_bytes(image_data["image"])
            continue

        pixmap = fitz.Pixmap(doc, item["xref"])
        try:
            if pixmap.alpha or pixmap.n > 4:
                pixmap = fitz.Pixmap(fitz.csRGB, pixmap)
            pixmap.save(str(output_path))
        finally:
            pixmap = None


def write_metadata(images_dir, metadata):
    metadata_path = images_dir / METADATA_FILE
    metadata_path.write_text(
        json.dumps(metadata, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )


def create_zip(source_dir, zip_path):
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as archive:
        for path in sorted(source_dir.iterdir()):
            archive.write(path, arcname=path.name)


def sanitize_pdf_basename(filename):
    stem = Path(filename).stem.strip() or "documento"
    stem = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "-", stem)
    stem = re.sub(r"\s+", " ", stem).strip(" .")
    return stem or "documento"


def write_session_info(workdir, original_filename):
    session_path = workdir / SESSION_FILE
    session_path.write_text(
        json.dumps({
            "original_filename": original_filename,
            "original_stem": sanitize_pdf_basename(original_filename),
            "modified_images": [],
        }, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def read_session_info(workdir):
    session_path = workdir / SESSION_FILE
    if not session_path.exists():
        return {}
    return json.loads(session_path.read_text(encoding="utf-8"))


def update_session_info(workdir, **updates):
    session_info = read_session_info(workdir)
    session_info.update(updates)
    (workdir / SESSION_FILE).write_text(
        json.dumps(session_info, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def mark_image_as_modified(workdir, filename):
    session_info = read_session_info(workdir)
    modified_images = set(session_info.get("modified_images", []))
    modified_images.add(filename)
    session_info["modified_images"] = sorted(modified_images)
    update_session_info(workdir, **session_info)


def build_updated_pdf_filename(workdir):
    session_info = read_session_info(workdir)
    stem = session_info.get("original_stem") or "documento"
    date_suffix = datetime.now().strftime("%Y-%m-%d")
    return f"{stem}-actualizado-{date_suffix}.pdf"


def build_compressed_pdf_filename(workdir):
    session_info = read_session_info(workdir)
    stem = session_info.get("original_stem") or "documento"
    date_suffix = datetime.now().strftime("%Y-%m-%d")
    return f"{stem}-comprimido-{date_suffix}.pdf"


def build_watermarked_pdf_filename(workdir, source="original"):
    session_info = read_session_info(workdir)
    stem = session_info.get("original_stem") or "documento"
    date_suffix = datetime.now().strftime("%Y-%m-%d")
    if source == "edited":
        return f"{stem}-editado-marca-agua-{date_suffix}.pdf"
    return f"{stem}-marca-agua-{date_suffix}.pdf"


def get_occurrence_rect(page_ref):
    rect_values = page_ref.get("rect")
    if rect_values and len(rect_values) == 4:
        return fitz.Rect(*rect_values)
    return None


def clamp(value, minimum, maximum):
    return max(minimum, min(value, maximum))


def parse_hex_color(color_value):
    color_value = (color_value or "#808080").strip()
    if not re.fullmatch(r"#[0-9a-fA-F]{6}", color_value):
        raise ValueError("Color de marca de agua no valido. Usa formato #RRGGBB.")
    return tuple(int(color_value[index:index + 2], 16) / 255 for index in (1, 3, 5))


def parse_watermark_config(form):
    words = (form.get("words") or "").strip()
    if not words:
        raise ValueError("La marca de agua necesita un texto.")

    config = {
        "words": words,
        "font_size": clamp(float(form.get("font_size", 36)), 8, 200),
        "color": parse_hex_color(form.get("color", "#808080")),
        "opacity": clamp(float(form.get("opacity", 0.18)), 0.05, 1),
        "rotate": clamp(float(form.get("rotate", -35)), -180, 180),
        "row": int(clamp(int(form.get("row", 3)), 1, 20)),
        "col": int(clamp(int(form.get("col", 2)), 1, 20)),
        "start_x": clamp(float(form.get("start_x", 90)), 0, 5000),
        "start_y": clamp(float(form.get("start_y", 140)), 0, 5000),
        "offset_x": clamp(float(form.get("offset_x", 220)), 20, 5000),
        "offset_y": clamp(float(form.get("offset_y", 180)), 20, 5000),
    }
    return config


def parse_bool_form(value, default=False):
    if value is None:
        return default
    return str(value).strip().lower() in {"1", "true", "yes", "on", "si"}


def save_pdf_document(doc, output_path, *, strip_metadata=False, linearize=False):
    if strip_metadata:
        doc.set_metadata({})

    save_kwargs = {
        "deflate": True,
        "garbage": 4,
        "clean": True,
    }

    if linearize:
        try:
            doc.save(output_path, linear=True, **save_kwargs)
            return
        except Exception as error:
            if "Linearisation is no longer supported" not in str(error):
                raise

    doc.save(output_path, **save_kwargs)


def apply_text_watermark(doc, config):
    for page in doc:
        page_width = page.rect.width
        page_height = page.rect.height
        for row_index in range(config["row"]):
            for col_index in range(config["col"]):
                x = min(config["start_x"] + (col_index * config["offset_x"]), page_width - 10)
                y = min(config["start_y"] + (row_index * config["offset_y"]), page_height - 10)
                point = fitz.Point(x, y)
                morph = None
                if config["rotate"]:
                    morph = (point, fitz.Matrix(config["rotate"]))
                page.insert_text(
                    point,
                    config["words"],
                    fontsize=config["font_size"],
                    fontname="helv",
                    color=config["color"],
                    fill_opacity=config["opacity"],
                    overlay=True,
                    morph=morph,
                )


def get_session_paths(session_id):
    workdir = DOWNLOADS_DIR / session_id
    images_dir = workdir / "imagenes"
    return workdir, images_dir


def get_session_download_dir(session_id):
    download_dir = DOWNLOADS_DIR / session_id
    download_dir.mkdir(parents=True, exist_ok=True)
    return download_dir


def cleanup_exported_pdfs(workdir, keep_name=None):
    for path in workdir.glob("*.pdf"):
        if keep_name and path.name == keep_name:
            continue
        path.unlink(missing_ok=True)


def get_safe_image_path(session_id, filename):
    if not filename or Path(filename).name != filename:
        raise ValueError("Nombre de archivo no valido")

    _workdir, images_dir = get_session_paths(session_id)
    image_path = images_dir / filename
    if not image_path.exists() or not image_path.is_file():
        raise FileNotFoundError(filename)

    return image_path, images_dir


def api_usage_response(endpoint, method="POST"):
    return jsonify({
        "ok": True,
        "endpoint": endpoint,
        "required_method": method,
        "message": f"Usa {method} para este endpoint.",
    })


def save_uploaded_pdf(file_storage):
    if file_storage is None:
        raise ValueError("No se envio un archivo PDF")
    if not file_storage.filename:
        raise ValueError("Archivo vacio")

    session_id = uuid.uuid4().hex[:12]
    workdir = DOWNLOADS_DIR / session_id
    workdir.mkdir(parents=True, exist_ok=True)
    write_session_info(workdir, file_storage.filename)
    pdf_path = workdir / SOURCE_FILE
    file_storage.save(str(pdf_path))

    if not pdf_path.exists() or pdf_path.stat().st_size == 0:
        remove_session(session_id)
        raise ValueError("El archivo PDF no se guardo correctamente")

    return session_id, workdir, pdf_path


def convert_pdf_to_docx(pdf_path, output_path):
    document = Document()
    normal_style = document.styles["Normal"]
    normal_style.font.name = "Calibri"
    normal_style.font.size = Pt(11)

    with fitz.open(pdf_path) as pdf:
        for page_index, page in enumerate(pdf):
            if page_index > 0:
                document.add_page_break()

            document.add_heading(f"Pagina {page_index + 1}", level=1)

            text_blocks = page.get_text("blocks", sort=True)
            text_found = False
            for block in text_blocks:
                text = block[4].strip()
                if not text:
                    continue

                lines = [line.strip() for line in text.splitlines() if line.strip()]
                if not lines:
                    continue

                document.add_paragraph("\n".join(lines))
                text_found = True

            if not text_found:
                document.add_paragraph("[Pagina sin texto extraible]")

            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            image_stream = BytesIO(pixmap.tobytes("png"))
            document.add_picture(image_stream, width=Inches(6.0))

    document.save(str(output_path))


def convert_pdf_to_xlsx(pdf_path, output_path):
    workbook = Workbook()
    default_sheet = workbook.active
    workbook.remove(default_sheet)

    with fitz.open(pdf_path) as pdf:
        for page_index, page in enumerate(pdf):
            sheet = workbook.create_sheet(title=f"Pagina {page_index + 1}")
            headers = ["Orden", "Texto", "X0", "Y0", "X1", "Y1"]
            sheet.append(headers)

            for cell in sheet[1]:
                cell.font = Font(bold=True)

            text_blocks = [
                block for block in page.get_text("blocks", sort=True)
                if len(block) >= 5 and block[4].strip()
            ]

            if not text_blocks:
                sheet.append([1, "[Pagina sin texto extraible]", "", "", "", ""])
            else:
                for order, block in enumerate(text_blocks, start=1):
                    x0, y0, x1, y1, text = block[:5]
                    normalized_text = "\n".join(
                        line.strip() for line in text.splitlines() if line.strip()
                    )
                    sheet.append([order, normalized_text, x0, y0, x1, y1])

            sheet.column_dimensions["A"].width = 10
            sheet.column_dimensions["B"].width = 80
            sheet.column_dimensions["C"].width = 12
            sheet.column_dimensions["D"].width = 12
            sheet.column_dimensions["E"].width = 12
            sheet.column_dimensions["F"].width = 12

    workbook.save(str(output_path))


def convert_pdf_to_pptx(pdf_path, output_path):
    if not PPTX_AVAILABLE:
        raise ValueError("PowerPoint no esta disponible. Instala python-pptx.")

    presentation = Presentation()
    presentation.slide_width = PptInches(13.333)
    presentation.slide_height = PptInches(7.5)
    blank_layout = presentation.slide_layouts[6]

    with fitz.open(pdf_path) as pdf:
        for page_index, page in enumerate(pdf, start=1):
            slide = presentation.slides.add_slide(blank_layout)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_stream = BytesIO(pixmap.tobytes("png"))
            slide.shapes.add_picture(
                image_stream,
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = f"Pagina {page_index}"
            page_text = page.get_text("text").strip()
            if page_text:
                notes_frame.add_paragraph().text = page_text[:4000]

    presentation.save(str(output_path))


def optimize_image_for_web(image_path, output_path, *, max_width=1600, quality=72, prefer_webp=True):
    if not PIL_AVAILABLE:
        raise ValueError("La optimizacion web de imagenes requiere Pillow.")

    max_width = max(320, min(int(max_width), 4000))
    quality = max(30, min(int(quality), 90))

    with Image.open(image_path) as image:
        image = image.convert("RGBA") if image.mode in {"LA", "RGBA", "P"} else image.convert("RGB")
        if image.width > max_width:
            ratio = max_width / image.width
            image = image.resize((max_width, int(image.height * ratio)), Image.LANCZOS)

        has_alpha = "A" in image.getbands()
        if prefer_webp:
            output_path = output_path.with_suffix(".webp")
            image.save(output_path, format="WEBP", quality=quality, method=6)
            return output_path

        output_path = output_path.with_suffix(".png" if has_alpha else ".jpg")
        save_format = "PNG" if has_alpha else "JPEG"
        save_image = image if has_alpha else image.convert("RGB")
        save_kwargs = {} if has_alpha else {"quality": quality, "optimize": True, "progressive": True}
        save_image.save(output_path, format=save_format, **save_kwargs)
        return output_path


def export_images_for_web(images_dir, output_zip_path, *, max_width=1600, quality=72, prefer_webp=True):
    if not images_dir.exists():
        raise ValueError("No hay imagenes extraidas para optimizar.")

    temp_dir = output_zip_path.with_suffix("")
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
    temp_dir.mkdir(parents=True, exist_ok=True)

    try:
        exported = 0
        for image_path in sorted(images_dir.iterdir()):
            if not image_path.is_file() or image_path.name in {METADATA_FILE}:
                continue
            output_path = temp_dir / image_path.stem
            optimize_image_for_web(
                image_path,
                output_path,
                max_width=max_width,
                quality=quality,
                prefer_webp=prefer_webp,
            )
            exported += 1

        if exported == 0:
            raise ValueError("No se encontraron imagenes para optimizar.")

        create_zip(temp_dir, output_zip_path)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def pixmap_to_image_bytes(pixmap, image_format, quality):
    image_format = image_format.lower()

    if image_format == "jpeg":
        return pixmap.tobytes("jpeg", jpg_quality=quality)

    if image_format == "webp":
        if not supports_webp_insert():
            raise RuntimeError("WebP no esta disponible para insercion PDF en este entorno")
        return pixmap.pil_tobytes(format="WEBP", quality=quality, method=6)

    raise ValueError(f"Formato de compresion no soportado: {image_format}")


def compress_pdf_rasterized(
    source_path,
    output_path,
    dpi=110,
    quality=55,
    image_format="jpeg",
    grayscale=False,
    strip_metadata=False,
    linearize=True,
):
    dpi = max(36, min(int(dpi), 180))
    quality = max(15, min(int(quality), 85))

    with fitz.open(source_path) as source_doc:
        compressed_doc = fitz.open()
        for page in source_doc:
            pixmap = page.get_pixmap(dpi=dpi, alpha=False)
            if grayscale and pixmap.colorspace.n != 1:
                pixmap = fitz.Pixmap(fitz.csGRAY, pixmap)
            image_bytes = pixmap_to_image_bytes(pixmap, image_format=image_format, quality=quality)
            new_page = compressed_doc.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(page.rect, stream=image_bytes)

        save_pdf_document(
            compressed_doc,
            output_path,
            strip_metadata=strip_metadata,
            linearize=linearize,
        )


def resolve_compression_candidates(level):
    presets = {
        "media": [
            {"dpi": 125, "quality": 68, "image_format": "webp"},
            {"dpi": 125, "quality": 68, "image_format": "jpeg"},
            {"dpi": 110, "quality": 60, "image_format": "webp"},
            {"dpi": 110, "quality": 60, "image_format": "jpeg"},
            {"dpi": 96, "quality": 50, "image_format": "webp"},
            {"dpi": 96, "quality": 50, "image_format": "jpeg"},
            {"dpi": 84, "quality": 42, "image_format": "webp"},
            {"dpi": 72, "quality": 35, "image_format": "jpeg"},
        ],
        "alta": [
            {"dpi": 105, "quality": 52, "image_format": "webp"},
            {"dpi": 105, "quality": 52, "image_format": "jpeg"},
            {"dpi": 96, "quality": 46, "image_format": "webp"},
            {"dpi": 96, "quality": 46, "image_format": "jpeg"},
            {"dpi": 84, "quality": 40, "image_format": "webp"},
            {"dpi": 72, "quality": 34, "image_format": "jpeg"},
        ],
        "maxima": [
            {"dpi": 96, "quality": 42, "image_format": "webp"},
            {"dpi": 96, "quality": 42, "image_format": "jpeg"},
            {"dpi": 84, "quality": 34, "image_format": "webp"},
            {"dpi": 84, "quality": 34, "image_format": "jpeg"},
            {"dpi": 72, "quality": 28, "image_format": "jpeg"},
            {"dpi": 60, "quality": 24, "image_format": "jpeg"},
            {"dpi": 50, "quality": 20, "image_format": "jpeg"},
            {"dpi": 42, "quality": 18, "image_format": "jpeg"},
            {"dpi": 36, "quality": 15, "image_format": "jpeg"},
        ],
    }
    candidates = presets.get(level, presets["media"])
    if supports_webp_insert():
        return candidates
    return [candidate for candidate in candidates if candidate["image_format"] != "webp"]


def compress_pdf_to_target(
    source_path,
    output_path,
    level,
    max_bytes=MAX_COMPRESSED_BYTES,
    grayscale=False,
    strip_metadata=False,
    linearize=True,
):
    candidates = resolve_compression_candidates(level)
    temp_paths = []
    best_path = None
    best_size = None
    target_met = False
    maximize_reduction = level == "maxima"

    try:
        for index, candidate in enumerate(candidates, start=1):
            temp_path = output_path.with_name(f"{output_path.stem}.tmp{index}{output_path.suffix}")
            temp_paths.append(temp_path)
            compress_pdf_rasterized(
                source_path,
                temp_path,
                dpi=candidate["dpi"],
                quality=candidate["quality"],
                image_format=candidate["image_format"],
                grayscale=grayscale,
                strip_metadata=strip_metadata,
                linearize=linearize,
            )
            temp_size = temp_path.stat().st_size
            if best_size is None or temp_size < best_size:
                best_size = temp_size
                best_path = temp_path
            if temp_size <= max_bytes:
                target_met = True
                if not maximize_reduction:
                    break

        if best_path is None:
            raise ValueError("No se pudo comprimir el PDF")

        shutil.copyfile(best_path, output_path)
        return {
            "compressed_size": best_size,
            "target_met": target_met,
            "max_bytes": max_bytes,
            "mode": "maxima" if maximize_reduction else "target",
        }
    finally:
        for temp_path in temp_paths:
            if temp_path.exists():
                temp_path.unlink(missing_ok=True)


def safe_extract_zip(zip_path, target_dir):
    target_root = target_dir.resolve()
    with zipfile.ZipFile(zip_path, "r") as archive:
        for member in archive.infolist():
            member_path = target_root / member.filename
            resolved_member = member_path.resolve()
            if target_root not in resolved_member.parents and resolved_member != target_root:
                raise ValueError("El ZIP contiene rutas no permitidas")

            if member.is_dir():
                resolved_member.mkdir(parents=True, exist_ok=True)
                continue

            resolved_member.parent.mkdir(parents=True, exist_ok=True)
            with archive.open(member, "r") as source, open(resolved_member, "wb") as target:
                shutil.copyfileobj(source, target)


@app.errorhandler(RequestEntityTooLarge)
def too_large(_error):
    return jsonify({"error": "El archivo excede 200 MB"}), 413


@app.errorhandler(413)
def too_large_generic(_error):
    return jsonify({"error": "El archivo excede el tamano maximo permitido"}), 413


@app.errorhandler(404)
def not_found(_error):
    if request.path.startswith("/api/"):
        return jsonify({"error": "Recurso no encontrado"}), 404
    return render_template("index.html"), 404


@app.errorhandler(405)
def method_not_allowed(_error):
    return jsonify({
        "error": f"Metodo no permitido para {request.path}. Usa POST."
    }), 405


@app.errorhandler(500)
def server_error(_error):
    app.logger.error(traceback.format_exc())
    return jsonify({"error": "Error interno del servidor"}), 500


@app.errorhandler(Exception)
def catch_all(error):
    if isinstance(error, HTTPException):
        return jsonify({"error": error.description}), error.code
    app.logger.error(traceback.format_exc())
    return jsonify({"error": "Error interno del servidor"}), 500


@app.after_request
def cors_headers(response):
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, OPTIONS"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type"
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    return response


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/pdf-info", methods=["POST"])
def api_pdf_info():
    ensure_upload_dir()
    session_id = None
    try:
        session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))

        with fitz.open(pdf_path) as doc:
            page_count = len(doc)
            first_page = doc[0] if page_count > 0 else None
            page_width = round(first_page.rect.width, 1) if first_page else 0
            page_height = round(first_page.rect.height, 1) if first_page else 0
            image_count = sum(1 for p in range(page_count) for _ in doc[p].get_images())

            meta = doc.metadata
            title = meta.get("title", "")
            author = meta.get("author", "")
            subject = meta.get("subject", "")
            creator = meta.get("creator", "")
            producer = meta.get("producer", "")
            creation = meta.get("creationDate", "")
            modified = meta.get("modDate", "")

            return jsonify({
                "session_id": session_id,
                "pages": page_count,
                "images": image_count,
                "page_width": page_width,
                "page_height": page_height,
                "page_size": f"{page_width} x {page_height} pts",
                "title": title,
                "author": author,
                "subject": subject,
                "creator": creator,
                "producer": producer,
                "created": creation,
                "modified": modified,
            })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        if session_id:
            remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/extract", methods=["GET", "POST", "OPTIONS"])
def api_extract():
    ensure_upload_dir()

    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/extract")

    try:
        session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
        images_dir = workdir / "imagenes"
        download_dir = get_session_download_dir(session_id)
        images_dir.mkdir(exist_ok=True)

        with fitz.open(pdf_path) as doc:
            metadata = build_image_metadata(doc)
            save_extracted_images(doc, images_dir, metadata)

        write_metadata(images_dir, metadata)
        create_zip(images_dir, download_dir / "imagenes.zip")

        return jsonify({
            "session_id": session_id,
            "total": len(metadata),
            "download_url": f"/api/download/{session_id}/imagenes.zip",
            "imagenes": metadata,
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/regenerate", methods=["GET", "POST", "OPTIONS"])
def api_regenerate():
    ensure_upload_dir()

    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/regenerate")

    session_id = request.form.get("session_id")
    if not session_id:
        return jsonify({"error": "Falta session_id"}), 400

    workdir = DOWNLOADS_DIR / session_id
    pdf_path = workdir / SOURCE_FILE
    images_dir = workdir / "imagenes"
    download_dir = get_session_download_dir(session_id)

    if not pdf_path.exists():
        return jsonify({"error": "Sesion no encontrada. Vuelve a extraer."}), 400

    meta_path = images_dir / METADATA_FILE
    if not meta_path.exists():
        return jsonify({"error": "No hay metadatos. Extrae primero."}), 400

    try:
        zip_file = request.files.get("zip")
        zip_modified_names = set()
        if zip_file and zip_file.filename:
            zip_path = workdir / "editadas.zip"
            zip_file.save(str(zip_path))
            with zipfile.ZipFile(zip_path, "r") as archive:
                zip_modified_names = {
                    Path(member.filename).name
                    for member in archive.infolist()
                    if not member.is_dir() and Path(member.filename).name
                }
            safe_extract_zip(zip_path, images_dir)

        metadata = json.loads(meta_path.read_text(encoding="utf-8"))
        session_info = read_session_info(workdir)
        modified_images = set(session_info.get("modified_images", []))
        modified_images.update(zip_modified_names)
        if zip_modified_names:
            update_session_info(workdir, modified_images=sorted(modified_images))

        replaced = 0
        output_filename = build_updated_pdf_filename(workdir)
        output_path = download_dir / output_filename

        if not modified_images:
            cleanup_exported_pdfs(download_dir, keep_name=output_filename)
            shutil.copyfile(pdf_path, output_path)
            return jsonify({
                "reemplazadas": 0,
                "download_url": f"/api/download/{session_id}/{output_filename}",
                "saved_name": output_filename,
                "saved_dir": str(download_dir),
                "saved_path": str(output_path),
            })

        with fitz.open(pdf_path) as doc:
            cleanup_exported_pdfs(download_dir, keep_name=output_filename)
            for item in metadata:
                if item["name"] not in modified_images:
                    continue

                image_path = images_dir / item["name"]
                if not image_path.exists():
                    continue

                try:
                    image_bytes = image_path.read_bytes()
                    target_page = item["pages"][0]["page"]
                    doc[target_page].delete_image(item["xref"])

                    for page_ref in item["pages"]:
                        rect = get_occurrence_rect(page_ref)
                        if rect is None:
                            continue
                        page = doc[page_ref["page"]]
                        page.insert_image(
                            rect,
                            stream=image_bytes,
                            keep_proportion=True,
                            overlay=True,
                        )
                    replaced += 1
                except Exception as image_error:
                    app.logger.warning("No se pudo reemplazar %s: %s", item["name"], image_error)

            doc.save(output_path, deflate=True, garbage=4, clean=True)

        return jsonify({
            "reemplazadas": replaced,
            "download_url": f"/api/download/{session_id}/{output_filename}",
            "saved_name": output_filename,
            "saved_dir": str(download_dir),
            "saved_path": str(output_path),
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/image/rotate", methods=["POST"])
def api_rotate_image():
    session_id = request.form.get("session_id")
    filename = request.form.get("filename")
    angle = request.form.get("angle", "90")

    if not session_id:
        return jsonify({"error": "Falta session_id"}), 400
    if not filename:
        return jsonify({"error": "Falta filename"}), 400

    try:
        angle = int(angle)
    except ValueError:
        return jsonify({"error": "Angulo debe ser un numero entero"}), 400

    try:
        image_path, images_dir = get_safe_image_path(session_id, filename)
        from PIL import Image
        with Image.open(str(image_path)) as img:
            rotated = img.rotate(-angle, expand=True)
            rotated.save(str(image_path))
        mark_image_as_modified(image_path.parent.parent, filename)
        create_zip(images_dir, get_session_download_dir(session_id) / "imagenes.zip")
        return jsonify({
            "ok": True,
            "image_url": f"/api/image/{session_id}/{filename}",
            "download_url": f"/api/download/{session_id}/imagenes.zip",
            "angle": angle,
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except FileNotFoundError:
        return jsonify({"error": "Imagen no encontrada"}), 404
    except Exception as error:
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/image/<session_id>/<filename>")
def api_image(session_id, filename):
    try:
        image_path, _images_dir = get_safe_image_path(session_id, filename)
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except FileNotFoundError:
        return jsonify({"error": "Imagen no encontrada"}), 404

    return send_file(str(image_path))


@app.route("/api/image/update", methods=["POST"])
def api_update_image():
    session_id = request.form.get("session_id")
    filename = request.form.get("filename")
    image_file = request.files.get("image")

    if not session_id:
        return jsonify({"error": "Falta session_id"}), 400
    if not filename:
        return jsonify({"error": "Falta filename"}), 400
    if image_file is None or not image_file.filename:
        return jsonify({"error": "Falta la imagen editada"}), 400

    try:
        image_path, images_dir = get_safe_image_path(session_id, filename)
        image_file.save(str(image_path))
        mark_image_as_modified(image_path.parent.parent, filename)
        create_zip(images_dir, get_session_download_dir(session_id) / "imagenes.zip")
        return jsonify({
            "ok": True,
            "image_url": f"/api/image/{session_id}/{filename}",
            "download_url": f"/api/download/{session_id}/imagenes.zip",
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except FileNotFoundError:
        return jsonify({"error": "Imagen no encontrada"}), 404
    except Exception as error:
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/convert-word", methods=["GET", "POST", "OPTIONS"])
def api_convert_word():
    ensure_upload_dir()

    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/convert-word")

    session_id = None
    try:
        session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
        output_path = get_session_download_dir(session_id) / "documento_convertido.docx"
        convert_pdf_to_docx(pdf_path, output_path)

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/documento_convertido.docx",
            "message": "Documento Word generado correctamente.",
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        if session_id:
            remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/convert-excel", methods=["GET", "POST", "OPTIONS"])
def api_convert_excel():
    ensure_upload_dir()

    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/convert-excel")

    session_id = None
    try:
        session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
        output_path = get_session_download_dir(session_id) / "documento_convertido.xlsx"
        convert_pdf_to_xlsx(pdf_path, output_path)

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/documento_convertido.xlsx",
            "message": "Archivo Excel generado correctamente.",
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        if session_id:
            remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/convert-pptx", methods=["GET", "POST", "OPTIONS"])
def api_convert_pptx():
    ensure_upload_dir()

    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/convert-pptx")

    try:
        session_id, _workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
        output_path = get_session_download_dir(session_id) / "documento_convertido.pptx"
        convert_pdf_to_pptx(pdf_path, output_path)

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/documento_convertido.pptx",
            "message": "Presentacion PowerPoint generada correctamente.",
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/export-web-images", methods=["GET", "POST", "OPTIONS"])
def api_export_web_images():
    ensure_upload_dir()

    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/export-web-images")

    session_id = request.form.get("session_id")
    if not session_id:
        return jsonify({"error": "Falta session_id. Extrae primero las imagenes."}), 400

    prefer_webp = parse_bool_form(request.form.get("prefer_webp"), default=True)
    max_width = int(request.form.get("max_width", 1600))
    quality = int(request.form.get("quality", 72))

    try:
        workdir, images_dir = get_session_paths(session_id)
        if not workdir.exists():
            return jsonify({"error": "Sesion no encontrada. Vuelve a extraer el PDF."}), 400

        output_name = "imagenes_web.zip"
        output_path = get_session_download_dir(session_id) / output_name
        export_images_for_web(
            images_dir,
            output_path,
            max_width=max_width,
            quality=quality,
            prefer_webp=prefer_webp,
        )

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/{output_name}",
            "message": "ZIP de imagenes optimizadas para web generado correctamente.",
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/compress", methods=["GET", "POST", "OPTIONS"])
def api_compress():
    ensure_upload_dir()

    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/compress")

    session_id = request.form.get("session_id")
    source = request.form.get("source", "original")
    level = request.form.get("level", "media")
    maximize_reduction = level == "maxima"
    grayscale = parse_bool_form(request.form.get("grayscale"))
    strip_metadata = parse_bool_form(request.form.get("strip_metadata"))
    linearize = parse_bool_form(request.form.get("linearize"), default=True)

    try:
        if request.files.get("pdf"):
            session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
            source_path = pdf_path
            download_dir = get_session_download_dir(session_id)
            output_name = "documento_comprimido.pdf"
        else:
            if not session_id:
                return jsonify({"error": "Falta session_id o un archivo PDF"}), 400

            workdir = DOWNLOADS_DIR / session_id
            download_dir = get_session_download_dir(session_id)
            if source == "edited":
                source_path = download_dir / build_updated_pdf_filename(workdir)
                output_name = build_updated_pdf_filename(workdir)
                if not source_path.exists():
                    return jsonify({"error": "Primero genera el PDF editado."}), 400
            else:
                source_path = workdir / SOURCE_FILE
                output_name = build_compressed_pdf_filename(workdir)
                if not source_path.exists():
                    return jsonify({"error": "Sesion no encontrada. Vuelve a subir el PDF."}), 400

        original_size = source_path.stat().st_size
        if original_size <= MAX_COMPRESSED_BYTES and not maximize_reduction:
            if source == "edited":
                effective_name = build_updated_pdf_filename(workdir)
                cleanup_exported_pdfs(download_dir, keep_name=effective_name)
                if source_path.name != effective_name:
                    shutil.copyfile(source_path, download_dir / effective_name)
            else:
                effective_name = build_compressed_pdf_filename(workdir)
                cleanup_exported_pdfs(download_dir, keep_name=effective_name)
                shutil.copyfile(source_path, download_dir / effective_name)

            return jsonify({
                "session_id": session_id,
                "download_url": f"/api/download/{session_id}/{effective_name}",
                "original_size": original_size,
                "compressed_size": original_size,
                "reduction_percent": 0,
                "target_size": MAX_COMPRESSED_BYTES,
                "message": "El PDF ya cumple el limite y no supera 1 MB.",
            })

        output_path = download_dir / output_name
        cleanup_exported_pdfs(download_dir, keep_name=output_name)
        compression_result = compress_pdf_to_target(
            source_path,
            output_path,
            level=level,
            max_bytes=MAX_COMPRESSED_BYTES,
            grayscale=grayscale,
            strip_metadata=strip_metadata,
            linearize=linearize,
        )
        compressed_size = compression_result["compressed_size"]

        if compressed_size >= original_size:
            output_path.unlink(missing_ok=True)
            return jsonify({
                "error": "No fue posible generar una version mas liviana sin superar el limite de 1 MB.",
            }), 400
        if not compression_result["target_met"]:
            output_path.unlink(missing_ok=True)
            return jsonify({
                "error": "No fue posible comprimir el PDF para que no supere 1 MB.",
            }), 400

        reduction = max(0, original_size - compressed_size)
        reduction_percent = round((reduction / original_size) * 100, 1) if original_size else 0

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/{output_name}",
            "original_size": original_size,
            "compressed_size": compressed_size,
            "reduction_percent": reduction_percent,
            "target_size": MAX_COMPRESSED_BYTES,
            "message": (
                "PDF comprimido al maximo posible."
                if maximize_reduction
                else "PDF comprimido correctamente. El archivo no supera 1 MB."
            ),
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        if request.files.get("pdf") and session_id:
            remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/watermark", methods=["GET", "POST", "OPTIONS"])
def api_watermark():
    ensure_upload_dir()

    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/watermark")

    session_id = request.form.get("session_id")
    source = request.form.get("source", "original")
    strip_metadata = parse_bool_form(request.form.get("strip_metadata"))

    try:
        config = parse_watermark_config(request.form)

        if request.files.get("pdf"):
            session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
            source_path = pdf_path
            download_dir = get_session_download_dir(session_id)
            output_name = "documento_marca_agua.pdf"
        else:
            if not session_id:
                return jsonify({"error": "Falta session_id o un archivo PDF"}), 400

            workdir = DOWNLOADS_DIR / session_id
            download_dir = get_session_download_dir(session_id)
            if source == "edited":
                source_path = download_dir / build_updated_pdf_filename(workdir)
                output_name = build_watermarked_pdf_filename(workdir, source="edited")
                if not source_path.exists():
                    return jsonify({"error": "Primero genera el PDF editado."}), 400
            else:
                source_path = workdir / SOURCE_FILE
                output_name = build_watermarked_pdf_filename(workdir, source="original")
                if not source_path.exists():
                    return jsonify({"error": "Sesion no encontrada. Vuelve a subir el PDF."}), 400

        output_path = download_dir / output_name
        output_path.unlink(missing_ok=True)

        with fitz.open(source_path) as doc:
            apply_text_watermark(doc, config)
            save_pdf_document(doc, output_path, strip_metadata=strip_metadata, linearize=True)

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/{output_name}",
            "saved_name": output_name,
            "saved_dir": str(download_dir),
            "saved_path": str(output_path),
            "message": "Marca de agua aplicada correctamente.",
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        if request.files.get("pdf") and session_id:
            remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error interno del servidor: {error}"}), 500


@app.route("/api/download/<session_id>/<filename>")
def api_download(session_id, filename):
    if not filename or Path(filename).name != filename:
        return jsonify({"error": "Nombre de archivo no valido"}), 400
    if Path(filename).suffix.lower() not in ALLOWED_EXPORT_SUFFIXES:
        return jsonify({"error": "Tipo de archivo no permitido para descarga"}), 400

    session_roots = [(DOWNLOADS_DIR / session_id).resolve()]
    filepath = None
    for root in session_roots:
        candidate = (root / filename).resolve()
        if root in candidate.parents and candidate.exists():
            filepath = candidate
            break

    if filepath is None:
        return jsonify({"error": "Archivo no encontrado"}), 404
    
    preview = request.args.get("preview", "false").lower() in {"1", "true", "yes", "on"}
    return send_file(str(filepath), as_attachment=not preview)


# ═══════════════════════════════════════════════════════════════════════════════
# MERGE / REORDER / ADD PAGE endpoints
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/merge", methods=["GET", "POST", "OPTIONS"])
def api_merge():
    ensure_upload_dir()
    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/merge")

    files = request.files.getlist("pdfs")
    if not files or len(files) < 2:
        return jsonify({"error": "Se necesitan al menos 2 PDFs"}), 400

    session_id = uuid.uuid4().hex[:12]
    workdir = get_session_download_dir(session_id)
    output_path = workdir / "documento_unido.pdf"

    try:
        merged_doc = fitz.open()
        for file in files:
            if not file.filename.lower().endswith(".pdf"):
                continue
            temp_path = workdir / f"temp_{file.filename}"
            file.save(str(temp_path))
            with fitz.open(temp_path) as doc:
                merged_doc.insert_pdf(doc)
            temp_path.unlink(missing_ok=True)

        merged_doc.save(output_path, deflate=True, garbage=4, clean=True)
        merged_doc.close()

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/documento_unido.pdf",
            "pages": len(merged_doc),
        })
    except Exception as error:
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error al unir PDFs: {error}"}), 500


@app.route("/api/reorder", methods=["GET", "POST", "OPTIONS"])
def api_reorder():
    ensure_upload_dir()
    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/reorder")

    session_id = None
    try:
        session_id = request.form.get("session_id")
        if session_id:
            workdir = DOWNLOADS_DIR / session_id
            pdf_path = workdir / SOURCE_FILE
            if not pdf_path.exists():
                return jsonify({"error": "Sesion no encontrada"}), 400
        else:
            session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
        
        order_str = request.form.get("order", "")
        if not order_str:
            return jsonify({"error": "Falta el orden de paginas"}), 400

        try:
            page_order = [int(x.strip()) - 1 for x in order_str.split(",")]
        except ValueError:
            return jsonify({"error": "Formato de orden invalido. Usa 1,3,2,4"}), 400

        download_dir = get_session_download_dir(session_id)
        output_path = download_dir / "documento_reordenado.pdf"

        with fitz.open(pdf_path) as doc:
            new_doc = fitz.open()
            for page_num in page_order:
                if 0 <= page_num < len(doc):
                    new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
            new_doc.save(output_path, deflate=True, garbage=4, clean=True)
            # Replace the source file with the reordered document so subsequent operations see the new order
            temp_path = workdir / "source_temp.bin"
            new_doc.save(temp_path, deflate=True, garbage=4, clean=True)
            new_doc.close()
        
        # Replace original file
        pdf_path.unlink(missing_ok=True)
        temp_path.rename(pdf_path)

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/documento_reordenado.pdf",
            "pages": len(page_order),
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        if session_id:
            remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error al reordenar: {error}"}), 500


@app.route("/api/add-page", methods=["GET", "POST", "OPTIONS"])
def api_add_page():
    ensure_upload_dir()
    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/add-page")

    session_id = None
    try:
        session_id = request.form.get("session_id")
        if session_id:
            workdir = DOWNLOADS_DIR / session_id
            pdf_path = workdir / SOURCE_FILE
            if not pdf_path.exists():
                return jsonify({"error": "Sesion no encontrada"}), 400
        else:
            session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
        
        position = request.form.get("position", "end")
        count = int(request.form.get("count", 1))
        page_number = request.form.get("page_number", "")

        download_dir = get_session_download_dir(session_id)
        output_path = download_dir / "documento_ampliado.pdf"

        with fitz.open(pdf_path) as doc:
            if position == "end":
                for _ in range(count):
                    doc.new_page(width=612, height=792)
            elif position == "start":
                for _ in range(count):
                    doc.insert_page(0, width=612, height=792)
            elif position == "after" and page_number:
                page_idx = int(page_number) - 1
                for _ in range(count):
                    doc.insert_page(page_idx + 1, width=612, height=792)
                    page_idx += 1
            elif position == "before" and page_number:
                page_idx = int(page_number) - 1
                for _ in range(count):
                    doc.insert_page(page_idx, width=612, height=792)
                    page_idx += 1

            total_pages = len(doc)
            doc.save(output_path, deflate=True, garbage=4, clean=True)

            # Replace the source file with the updated document so subsequent operations see all pages
            # Save to temp file first, then replace original (can't save incremental to same file)
            temp_path = workdir / "source_temp.bin"
            doc.save(temp_path, deflate=True, garbage=4, clean=True)
        
        # Close doc and replace original file
        pdf_path.unlink(missing_ok=True)
        temp_path.rename(pdf_path)

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/documento_ampliado.pdf",
            "pages": total_pages,
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        if session_id:
            remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error al agregar paginas: {error}"}), 500


@app.route("/api/thumbnails", methods=["POST", "GET"])
def api_thumbnails():
    app.logger.info("api_thumbnails called")
    session_id = request.form.get("session_id")
    if not session_id:
        return jsonify({"error": "Falta session_id"}), 400

    try:
        workdir = DOWNLOADS_DIR / session_id
        pdf_path = workdir / SOURCE_FILE
        app.logger.info(f"Looking for PDF at: {pdf_path}, exists: {pdf_path.exists()}")
        if not pdf_path.exists():
            return jsonify({"error": "PDF no encontrado"}), 404

        thumbs_dir = workdir / "thumbnails"
        thumbs_dir.mkdir(exist_ok=True)

        thumbnail_urls = []
        with fitz.open(pdf_path) as doc:
            for i, page in enumerate(doc):
                thumb_path = thumbs_dir / f"page_{i+1:03d}.png"
                if not thumb_path.exists():
                    pix = page.get_pixmap(matrix=fitz.Matrix(0.4, 0.4))
                    pix.save(str(thumb_path))
                thumbnail_urls.append(f"/api/thumbnail/{session_id}/page_{i+1:03d}.png")

        return jsonify({
            "session_id": session_id,
            "total_pages": len(thumbnail_urls),
            "urls": thumbnail_urls,
        })
    except Exception as error:
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error al generar thumbnails: {error}"}), 500


@app.route("/api/thumbnail/<session_id>/<filename>")
def api_thumbnail_image(session_id, filename):
    try:
        thumb_path = (DOWNLOADS_DIR / session_id / "thumbnails" / filename).resolve()
        if DOWNLOADS_DIR / session_id not in thumb_path.parents:
            raise ValueError("Ruta no valida")
        if not thumb_path.exists():
            return jsonify({"error": "Thumbnail no encontrado"}), 404
        return send_file(str(thumb_path))
    except Exception as error:
        app.logger.error(traceback.format_exc())
        return jsonify({"error": str(error)}), 400


@app.route("/api/foliar", methods=["GET", "POST", "OPTIONS"])
def api_foliar():
    ensure_upload_dir()
    if request.method in {"GET", "OPTIONS"}:
        return api_usage_response("/api/foliar")

    session_id = None
    try:
        session_id, workdir, pdf_path = save_uploaded_pdf(request.files.get("pdf"))
        start = int(request.form.get("start", 1))
        position = request.form.get("position", "bottom")
        prefix = request.form.get("prefix", "")
        font_size = int(request.form.get("font_size", 12))
        color = request.form.get("color", "#000000")
        format = request.form.get("format", "{n}")

        download_dir = get_session_download_dir(session_id)
        output_path = download_dir / "documento_foliado.pdf"

        with fitz.open(pdf_path) as doc:
            total_pages = len(doc)
            for i, page in enumerate(doc):
                page_num = start + i
                text = format.replace("{n}", str(page_num)).replace("{p}", str(total_pages))
                if prefix:
                    text = f"{prefix}{text}"

                rect = page.rect
                if position == "bottom":
                    x = rect.width / 2
                    y = rect.height - 40
                elif position == "top":
                    x = rect.width / 2
                    y = 30
                elif position == "center":
                    x = rect.width / 2
                    y = rect.height / 2
                elif position == "right":
                    x = rect.width - 40
                    y = rect.height / 2
                elif position == "left":
                    x = 40
                    y = rect.height / 2
                else:
                    x = rect.width / 2
                    y = rect.height - 40

                page.insert_text(
                    (x, y),
                    text,
                    fontsize=font_size,
                    color=parse_hex_color(color),
                    fontname="helv",
                )

            doc.save(output_path, deflate=True, garbage=4, clean=True)

        return jsonify({
            "session_id": session_id,
            "download_url": f"/api/download/{session_id}/documento_foliado.pdf",
            "pages": total_pages,
            "message": "Foliado aplicado correctamente.",
        })
    except ValueError as error:
        return jsonify({"error": str(error)}), 400
    except Exception as error:
        if session_id:
            remove_session(session_id)
        app.logger.error(traceback.format_exc())
        return jsonify({"error": f"Error al foliar: {error}"}), 500


@app.route("/api/test", methods=["GET", "POST"])
def api_test():
    return jsonify({"message": "Test OK"})


if __name__ == "__main__":
    ensure_upload_dir()
    port = int(os.environ.get("PORT", 5000))
    app.run(debug=False, host="0.0.0.0", port=port)
